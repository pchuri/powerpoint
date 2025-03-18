from pptx.chart import chart
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from typing import Literal, Union, List, Dict, Any

class ChartManager:
    def __init__(self):
        self.name = "Chart Manager"

    def determine_chart_type(self, data: Dict[str, Any]) -> tuple[XL_CHART_TYPE, str]:
        """
        Analyze the data structure and determine the most appropriate chart type.
        Uses heuristics based on data patterns to select the best visualization.
        
        Args:
            data: Dictionary containing chart data with keys for 'series', 'categories', etc.
            
        Returns:
            tuple: (PowerPoint chart type enum, chart_format string)
        """
        # Get basic data properties
        series_count = len(data["series"])
        categories = data.get("categories", [])
        category_count = len(categories) if categories else 0
        
        # Validate series data exists
        if not data["series"] or not all(s.get("values") for s in data["series"]):
            # Default to column chart if data is missing or incomplete
            return XL_CHART_TYPE.COLUMN_CLUSTERED, "category"

        # Check for XY scatter data (coordinates)
        # We look for values that are [x,y] pairs
        is_xy_data = False
        for series in data["series"]:
            values = series.get("values", [])
            if values:
                # Check if the first value is a coordinate pair
                first_value = values[0]
                is_xy_data = isinstance(first_value, (list, tuple)) and len(first_value) == 2
                
                # Verify all values are valid coordinate pairs
                if is_xy_data:
                    is_xy_data = all(
                        isinstance(v, (list, tuple)) and len(v) == 2 and 
                        all(isinstance(coord, (int, float)) for coord in v)
                        for v in values if v is not None
                    )
                break

        if is_xy_data:
            return XL_CHART_TYPE.XY_SCATTER, "xy"

        # Check for pie chart conditions:
        # 1. Single series
        # 2. Few categories (<=8 for readability)
        # 3. Values sum to approximately 100 (suggesting percentages)
        if series_count == 1 and categories and len(categories) <= 8:
            values = data["series"][0].get("values", [])
            try:
                # Filter out None or non-numeric values
                numeric_values = [float(v) for v in values if v is not None]
                if numeric_values:
                    total = sum(numeric_values)
                    
                    # Check if values are percentages (sum ≈ 100)
                    if 95 <= total <= 105:
                        return XL_CHART_TYPE.PIE, "category"
                    
                    # Check if all values are similar magnitude (good for pie)
                    if numeric_values:
                        max_val = max(numeric_values)
                        min_val = min(numeric_values)
                        # If largest value is less than 10x smallest, pie is still readable
                        if min_val > 0 and max_val / min_val < 10:
                            return XL_CHART_TYPE.PIE, "category"
            except (TypeError, ValueError):
                pass  # Non-numeric data, continue to other chart types

        # Check for time series patterns in category names
        time_related_terms = ["date", "time", "year", "month", "day", "quarter", "q1", "q2", "q3", "q4", 
                             "jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec"]
        
        has_time_categories = False
        if categories:
            has_time_categories = any(
                isinstance(cat, (str, int)) and 
                any(term in str(cat).lower() for term in time_related_terms)
                for cat in categories
            )
            
            # For time series data, use line chart
            if has_time_categories:
                # For multiple series over time, use line chart
                return XL_CHART_TYPE.LINE, "category"

        # Data size-based decisions
        if category_count > 0:
            # Many categories with few series: use column chart
            if category_count > 10 and series_count == 1:
                return XL_CHART_TYPE.COLUMN_CLUSTERED, "category"
                
            # Many series with few categories: use bar chart (better label readability)
            if series_count > 3 and category_count <= 10:
                return XL_CHART_TYPE.BAR_CLUSTERED, "category"
                
            # Many categories with multiple series: use line chart (less visual clutter)
            if category_count > 10 and series_count > 1:
                return XL_CHART_TYPE.LINE, "category"
        
        # Default recommendations based on series count
        if series_count > 1:
            return XL_CHART_TYPE.BAR_CLUSTERED, "category"
        else:
            return XL_CHART_TYPE.COLUMN_CLUSTERED, "category"


    def add_chart_to_slide(self, slide, chart_type: XL_CHART_TYPE, data: Dict[str, Any],
                           chart_format: str = "category") -> chart:
        """Add a chart to the slide with the specified data."""
        # Position chart in the middle of the slide with margins
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)

        if chart_format == "category":
            chart_data = CategoryChartData()
            chart_data.categories = data.get("categories", [])

            # Add each series
            for series in data["series"]:
                chart_data.add_series(series["name"], series["values"])

        elif chart_format == "xy":
            chart_data = XyChartData()

            # Add each series
            for series in data["series"]:
                series_data = chart_data.add_series(series["name"])
                for x, y in series["values"]:
                    series_data.add_data_point(x, y)

        # Add and configure the chart
        graphic_frame = slide.shapes.add_chart(
            chart_type, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart

        # Basic formatting
        chart.has_legend = True
        if len(data["series"]) > 1:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM

        # Add axis titles if provided
        if "x_axis" in data:
            chart.category_axis.axis_title.text_frame.text = data["x_axis"]
        if "y_axis" in data:
            chart.value_axis.axis_title.text_frame.text = data["y_axis"]

        return chart